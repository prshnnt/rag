"""Document loaders. Return list[dict] with keys: text, page (optional)."""
from __future__ import annotations

import json
from pathlib import Path


def _ext(path: str | Path) -> str:
    return Path(path).suffix.lower().lstrip(".")


def load_txt(path: str | Path) -> list[dict]:
    p = Path(path)
    return [{"text": p.read_text(encoding="utf-8", errors="replace")}]


def load_json(path: str | Path) -> list[dict]:
    data = json.loads(Path(path).read_text(encoding="utf-8"))
    return [{"text": json.dumps(data, indent=2)}]


def load_csv(path: str | Path) -> list[dict]:
    try:
        import csv
    except ImportError as e:
        raise RuntimeError("csv unavailable") from e
    rows: list[dict] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.DictReader(f):
            rows.append({"text": json.dumps(row, ensure_ascii=False)})
    return rows


def load_pdf(path: str | Path) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise RuntimeError("pypdf required: uv add pypdf") from e
    reader = PdfReader(str(path))
    out: list[dict] = []
    for i, page in enumerate(reader.pages):
        out.append({"text": page.extract_text() or "", "page": i + 1})
    return out


def load_docx(path: str | Path) -> list[dict]:
    try:
        from docx import Document  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("docx not installed: uv add python-docx") from e
    doc = Document(str(path))
    return [{"text": "\n".join(p.text for p in doc.paragraphs)}]


def load_xlsx(path: str | Path) -> list[dict]:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("openpyxl not installed: uv add openpyxl") from e
    wb = load_workbook(str(path), read_only=True, data_only=True)
    out: list[dict] = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = [str(c) if c is not None else "" for c in rows[0]]
        for r in rows[1:]:
            cells = [str(c) if c is not None else "" for c in r]
            out.append({"text": "\t".join(f"{h}: {v}" for h, v in zip(header, cells))})
    wb.close()
    return out


def load_pptx(path: str | Path) -> list[dict]:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as e:
        raise RuntimeError("python-pptx not installed: uv add python-pptx") from e
    prs = Presentation(str(path))
    out: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        texts.append(t)
        out.append({"text": "\n".join(texts), "page": i})
    return out


_LOADERS = {
    "txt": load_txt,
    "json": load_json,
    "csv": load_csv,
    "pdf": load_pdf,
    "docx": load_docx,
    "xlsx": load_xlsx,
    "pptx": load_pptx,
}


def get_json_schema(obj) -> dict:
    """Pydantic v2 helper, kept for parity with prior stub."""
    if hasattr(obj, "model_json_schema"):
        return obj.model_json_schema()
    raise TypeError(f"No JSON schema for {type(obj).__name__}")


def load(path: str | Path) -> list[dict]:
    """Dispatch by extension. Unknown ext -> RuntimeError."""
    ext = _ext(path)
    if ext not in _LOADERS:
        raise RuntimeError(f"Unsupported extension: .{ext}")
    return _LOADERS[ext](path)
